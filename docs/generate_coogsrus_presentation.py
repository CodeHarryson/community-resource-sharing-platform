from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os
import time

OUT = "docs/CoogsRUs_Presentation.pptx"
LOGO_PATH = "client/public/uh-logo.png"  # optional

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

def add_title_slide(title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_tf = slide.shapes.title.text_frame
    title_tf.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle
    if os.path.exists(LOGO_PATH):
        left = Inches(11.7); top = Inches(0.2); height = Inches(1.0)
        slide.shapes.add_picture(LOGO_PATH, left, top, height=height)

def add_bullets_slide(title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        p = body.add_paragraph() if i>0 else body.paragraphs[0]
        p.text = b
        p.level = 0
        p.font.size = Pt(18)


def add_slide_with_notes(title, bullets, notes):
    """Create a bullet slide and attach speaker notes."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        p = body.add_paragraph() if i>0 else body.paragraphs[0]
        p.text = b
        p.level = 0
        p.font.size = Pt(18)

    # speaker notes
    notes_slide = slide.notes_slide
    notes_tf = notes_slide.notes_text_frame
    notes_tf.clear()
    notes_tf.text = notes

def add_two_column_slide(title, left_lines, right_lines):
    # Create a blank slide and draw two text boxes for left/right columns.
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # set the title if present in the layout, otherwise add one
    try:
        slide.shapes.title.text = title
    except Exception:
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.0), Inches(0.6))
        title_tf = title_box.text_frame
        title_tf.text = title

    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(6.0), Inches(5.5))
    left_tf = left_box.text_frame
    left_tf.word_wrap = True
    left_tf.clear()
    for i, l in enumerate(left_lines):
        p = left_tf.add_paragraph() if i>0 else left_tf.paragraphs[0]
        p.text = l
        p.level = 0
        p.font.size = Pt(16)

    right_box = slide.shapes.add_textbox(Inches(6.7), Inches(1.0), Inches(6.0), Inches(5.5))
    right_tf = right_box.text_frame
    right_tf.word_wrap = True
    right_tf.clear()
    for i, l in enumerate(right_lines):
        p = right_tf.add_paragraph() if i>0 else right_tf.paragraphs[0]
        p.text = l
        p.level = 0
        p.font.size = Pt(16)

# Slides
add_title_slide("Coogs R Us — Project Walkthrough", "Community resource app (University of Houston)")

add_bullets_slide("Presentation Agenda", [
    "Project overview & goals",
    "Tech stack and architecture",
    "Setup & run (Docker Compose)",
    "Demo: core features",
    "DB seeding & images",
    "Notifications & endpoints",
    "Troubleshooting & contribution"
])

add_bullets_slide("Project Overview", [
    "Community resource sharing platform for UH students",
    "Users can list resources, request items, and get notifications",
    "Focused for class / student org demo"
])

add_two_column_slide("Tech Stack", 
    ["Frontend: React (CRA), MUI, Axios",
     "Backend: Node.js (Express), JWT, bcrypt",
     "Database: PostgreSQL (init.sql seed)",
     "Deployment: Docker + docker-compose, nginx (client)"],
    ["Storage: static images in client/public/images",
     "Auth: JWT middleware",
     "Notifications: DB + API endpoints",
     "Build: client static served by nginx"]
)

add_bullets_slide("Local Setup (Prereqs)", [
    "Install Docker & Docker Compose",
    "Install Python (optional, for generating slides)",
    "Ensure ports 3000 (client) and 5000 (server) are free"
])

add_bullets_slide("Start the App (Commands)", [
    "From repo root:",
    "docker-compose down -v  # (optional) reset DB",
 "docker-compose up -d --build",
    "Visit: http://localhost:3000 (frontend) and http://localhost:5000 (health)"
])

add_bullets_slide("Demo Flow: Quick Actions", [
    "Register a user or use seeded accounts",
    "Create a resource (title, category, optional image)",
    "Request a resource as another user",
    "Owner approves/denies — notifications are generated"
])

add_bullets_slide("DB Seeding & Reset", [
    "Initial seeds live in server/init.sql (users, resources, notifications)",
    "Init script runs only on empty Postgres volume",
    "To re-seed: docker-compose down -v && docker-compose up --build",
    "Or run server/scripts/apply_seeds.js to inject seeds into existing DB"
])

add_bullets_slide("Images & Uploads", [
    "Static images: place under client/public/images/",
    "Image path served at /images/<filename>",
    "After adding files: docker-compose build client && docker-compose up -d --no-deps --build client",
    "Option: implement image upload endpoint to avoid rebuilds"
])

add_bullets_slide("Notifications (UI & API)", [
    "GET /api/notifications — list for authenticated user",
    "PATCH /api/notifications/:id — mark read",
    "POST created when requests are made or status changes",
    "UI: NotificationBell shows unread count and popover"
])

add_bullets_slide("Key Endpoints (summary)", [
    "Auth: POST /api/auth/register, /login",
    "Resources: GET/POST /api/resources",
    "Requests: POST /api/requests, PATCH /api/requests/:id, GET /api/requests",
    "Notifications: GET/POST/PATCH/DELETE /api/notifications"
])

add_bullets_slide("Troubleshooting Tips", [
    "Blank frontend? Check browser console for JS errors",
    "Server crashes? docker-compose logs --tail=200 server",
    "DB seeds not applied? You likely need to remove the DB volume",
    "Images not showing? Ensure client/public/images/<file> exists and client rebuilt"
])

add_bullets_slide("How to Contribute / Extend", [
    "Add image upload endpoint + storage (S3 or local)",
    "Add tests for API endpoints (jest/supertest)",
    "Add migrations (node-pg-migrate) for schema evolution",
    "Polish UI, accessibility, and mobile responsiveness"
])

add_bullets_slide("Appendix / Resources", [
    "Repo: (your local repo path)",
    "Commands cheat-sheet (in docs folder)",
    "Contact: project maintainer / student org"
])

# -------------------------
# Step-by-step lab slides
# -------------------------
add_bullets_slide("Prerequisites", [
    "Git, Docker Desktop (with Compose)",
    "Optional: Node.js & npm (for local dev), Python (for docs)",
    "Ports: ensure 3000 and 5000 are free"
])

add_bullets_slide("Clone & Inspect", [
    "git clone https://github.com/CodeHarryson/community-resource-sharing-platform.git",
    "cd community-resource-sharing-platform",
    "Explore: client/, server/, docker-compose.yml"
])

add_bullets_slide("Start with Docker (recommended)", [
    "From repo root:",
    "docker-compose up -d --build",
    "Open: http://localhost:3000 (frontend), http://localhost:5000 (api health)"
])

add_bullets_slide("Reset DB & Re-seed (if needed)", [
    "Warning: this deletes DB data",
    "docker-compose down -v",
    "docker-compose up -d --build (init.sql runs on fresh DB)",
    "Or: docker-compose exec server node scripts/apply_seeds.js to add seeds to an existing DB"
])

add_bullets_slide("Add images & make them live", [
    "Place static images in client/public/images/",
    "Prefer hyphenated filenames (no spaces), e.g. kids-bicycle.jpg",
    "After adding: docker-compose build client && docker-compose up -d --no-deps --build client",
    "Alternative: implement image upload endpoint to avoid rebuilds"
])

add_bullets_slide("Common dev workflows", [
    "Run client dev server for fast frontend edits: cd client && npm install && npm start",
    "Run server locally for quick backend testing: cd server && npm install && npm run dev",
    "Use docker-compose build <service> to rebuild a single service"
])

add_bullets_slide("API quick examples", [
    "Register: POST /api/auth/register {name,email,password}",
    "Login: POST /api/auth/login {email,password} => token",
    "Create resource: POST /api/resources (Authorization header required)",
    "Request: POST /api/requests {resource_id}",
    "Notifications: GET /api/notifications (auth)"
])

add_bullets_slide("Troubleshooting checklist", [
    "Blank UI: open browser DevTools console for errors",
    "Server errors: docker-compose logs --tail=200 server",
    "Images missing: check client/public/images and rebuild client",
    "Seeds missing: recreate DB volume or run seed helper"
])

add_bullets_slide("Suggested student lab tasks", [
    "Implement image upload (multer) + attach image_url to resources",
    "Add tests for key API endpoints (Jest + Supertest)",
    "Add migrations and a CI step to run tests on PRs"
])

# Detailed step-by-step slides with speaker notes
add_slide_with_notes(
    "Step 1 — Clone & Inspect",
    ["git clone https://github.com/CodeHarryson/community-resource-sharing-platform.git", "cd community-resource-sharing-platform", "Open folders: client/, server/, docker-compose.yml"],
    "Walk students through the repository layout. Explain where frontend, backend, and DB live. Ask them to open key files: client/src/App.js, server/index.js, server/init.sql."
)

add_slide_with_notes(
    "Step 2 — Start the App (Docker)",
    ["docker-compose up -d --build", "docker-compose ps", "Open http://localhost:3000 and http://localhost:5000"],
    "Explain Docker Compose: three services (client, server, db). Demonstrate docker-compose ps and docker-compose logs --tail=200 server. If services aren't healthy, show how to inspect logs and fix issues."
)

add_slide_with_notes(
    "Step 3 — Register, Login, Fetch API",
    ["POST /api/auth/register", "POST /api/auth/login => get JWT token", "GET /api/resources (public)"],
    "Show how to use Postman or curl. Emphasize Authorization header format: 'Bearer <token>'. Let students try register/login and inspect JSON responses."
)

add_slide_with_notes(
    "Step 4 — Create Resource & Request Flow",
    ["Create resource (title, category, optional image)", "Login as another user and POST /api/requests", "Owner approves -> notifications generated"],
    "Walk through UI flows in the app: Create Resource, Dashboard (owner view), and Request actions. Explain DB tables (resources, requests, notifications) and how the server updates credits on approve."
)

add_slide_with_notes(
    "Step 5 — Images & Client Build",
    ["Add image: client/public/images/<filename>", "Rebuild client: docker-compose build client", "Restart client: docker-compose up -d --no-deps --build client"],
    "Explain why the client is a static build served by nginx. Explain choice of public/ folder. Demonstrate adding an image and making it live. Discuss why an upload endpoint avoids rebuilding."
)

add_slide_with_notes(
    "Step 6 — Re-seeding DB",
    ["To reset DB (destructive): docker-compose down -v", "docker-compose up -d --build", "Or run seed helper: docker-compose exec server node scripts/apply_seeds.js"],
    "Warn about removing volumes in a multi-user setting. Explain init.sql runs only when container initializes an empty DB data directory. Show the seed helper as a non-destructive alternative."
)

add_slide_with_notes(
    "Step 7 — Dev Workflow (fast iteration)",
    ["Run client dev server: cd client && npm start", "Run server locally: cd server && npm run dev"],
    "Recommend running the React dev server and server locally during development so students can edit and get hot reload without rebuilding Docker images. Explain cross-origin: client dev server calls server at http://localhost:5000 (CORS enabled)."
)

add_slide_with_notes(
    "Instructor Notes & Exercises",
    ["Implement image upload (multer) + integrate with CreateResource", "Add tests with Jest + Supertest", "Add migrations (node-pg-migrate)"],
    "Provide starter code pointers: where to add multer middleware and how to serve uploads, example tests to write, and how to add migrations. Suggest pairing students and defining acceptance criteria."
)

# Interactive exercises / questions for students
add_slide_with_notes(
    "Exercise 1 — Where to implement image upload?",
    ["Q: Which files should you change to add an upload endpoint?", "A: server/index.js (static serve), server/routes/uploads.js (new), client/CreateResource.js (file input + upload)"],
    "Ask students to map the change to files. Hint: Add multer to server dependencies, create POST /api/uploads returning the public URL, and update CreateResource to POST file first."
)

add_slide_with_notes(
    "Exercise 2 — Debug a blank page",
    ["Q: The app shows a blank page. Which file likely caused it?", "Choices: client/src/api.js, client/src/pages/Dashboard.js, client/src/components/ItemCard.js"],
    "Explain how to triage: open DevTools console, look for stack traces and file/line numbers. Hint: common causes are missing imports, wrong prop casing, or thrown exceptions during render. Have students reproduce and fix an intentional bug."
)

add_slide_with_notes(
    "Exercise 3 — API test task",
    ["Q: Write a Jest + Supertest test for POST /api/auth/register", "What should the test assert? (status, returned user shape, token present)"],
    "Guide students: create server/tests/auth.test.js, boot the server or import the app, and run the test. Show a minimal example structure and the assertions to include."
)

add_slide_with_notes(
    "Group Challenge — Notifications CRUD",
    ["Extend notifications: add 'mark all read' endpoint", "Add client button to call it and update UI state"],
    "Acceptance: endpoint PATCH /api/notifications/mark-all should mark all notifications for the authenticated user as read. Client button should call it and update unread badge. Encourage groups to split backend/frontend work."
)

# -------------------------
# Hands-on examples for backend/frontend/docker
# -------------------------
add_slide_with_notes(
    "Hands-on: Backend — Upload endpoint",
    [
        "Add a new route: server/routes/uploads.js using multer",
        "Store uploads in server/uploads and return { filename, url }",
        "Expose uploads via express.static('/uploads') in server/index.js",
        "Update POST /api/resources to accept image_filename and save it"
    ],
    "Step-by-step:\n1) Create branch: git checkout -b feat/upload-endpoint\n2) npm install multer (in server)\n3) Implement server/routes/uploads.js (multer.single('file')) to save and return URL\n4) In server/index.js add: app.use('/uploads', express.static(path.join(__dirname, 'uploads')));\n5) Test with curl: curl -F \"file=@./kids-bicycle.jpg\" http://localhost:5000/api/uploads"
)

add_slide_with_notes(
    "Hands-on: Frontend — File upload & form",
    [
        "Add file input to CreateResource form (client/src/pages/CreateResource.js)",
        "On submit: POST file to /api/uploads to get filename/url, then POST resource with image_filename",
        "Show preview of selected image and handle upload errors"
    ],
    "Step-by-step:\n1) Add <input type=\"file\" /> and local state for file\n2) Use API.post('/uploads', formData) to upload (set content-type multipart/form-data)\n3) After upload success, include image_filename in POST /api/resources body\n4) Run client dev server: cd client && npm start and test the flow"
)

add_slide_with_notes(
    "Hands-on: Docker — Persist uploads & deploy",
    [
        "Mount ./server/uploads as a volume in docker-compose so uploads persist",
        "Ensure server serves /uploads static route; rebuild server image when code changes",
        "Option: serve uploads from a CDN or S3 for production"
    ],
    "Commands:\n1) Add to docker-compose.yml under server: volumes: - ./server/uploads:/usr/src/app/uploads\n2) docker-compose build server && docker-compose up -d server\n3) To test: upload a file and visit http://localhost:5000/uploads/<filename>"
)

add_bullets_slide("Paths Covered", [
    "Local dev: run client (npm start) + server (npm run dev)",
    "Docker Compose: build & run entire stack (docker-compose up -d --build)",
    "Hybrid: run client locally for fast UI dev while server runs in Docker"
])

add_slide_with_notes(
    "Wrap-up: deliverables for students",
    ["Working app: register/login, create resource, request/approve, notifications", "One extension: image upload OR tests OR mark-all-read"],
    "Ask each group to open a PR with a short description, link to their branch, and a demo video or screenshots. Provide grading rubric: functionality (60%), code quality/tests (25%), demo (15%)."
)

# Save presentation
os.makedirs(os.path.dirname(OUT), exist_ok=True)
# On Windows we can hit PermissionError if the file is open; try to remove an existing file first
if os.path.exists(OUT):
    try:
        os.remove(OUT)
    except Exception as e:
        # Warn but continue; we'll attempt an atomic replace via a temp file below
        print(f"Warning: could not remove existing presentation {OUT}: {e}")

# Save to a temporary file and atomically replace the target to reduce race/permission issues
tmp_out = OUT + ".tmp"
try:
    prs.save(tmp_out)
    # os.replace is atomic on most platforms and will overwrite OUT if possible
    try:
        os.replace(tmp_out, OUT)
        print(f"Presentation written to: {OUT}")
    except PermissionError as e:
        # The target file is locked (commonly on Windows). Fall back to a timestamped file.
        ts = int(time.time())
        alt_out = OUT.replace('.pptx', f'.{ts}.pptx')
        try:
            os.replace(tmp_out, alt_out)
            print(f"Could not overwrite {OUT} (file locked). Wrote presentation to: {alt_out}")
        except Exception:
            # If even moving to alt_out fails, cleanup and re-raise
            if os.path.exists(tmp_out):
                try: os.remove(tmp_out)
                except Exception: pass
            raise
except Exception:
    # Cleanup temp file if something went wrong
    try:
        if os.path.exists(tmp_out):
            os.remove(tmp_out)
    except Exception:
        pass
    # Re-raise so the caller sees the error
    raise
